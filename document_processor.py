import os
import re
import docx2txt
import pandas as pd
from pypdf import PdfReader
from pathlib import Path
import gc  # Added explicit import for garbage collection

class DocumentProcessor:
    def __init__(self, docs_directory='downloaded_files', output_directory='processed_docs', max_file_size_mb=50):
        self.docs_directory = docs_directory
        self.output_directory = output_directory
        self.max_file_size_bytes = max_file_size_mb * 1024 * 1024  # Convert MB to bytes
        self.supported_extensions = ['.pdf', '.docx', '.doc', '.pptx', '.csv', '.xlsx', '.xls']
        os.makedirs(output_directory, exist_ok=True)
        
    def extract_text_from_pdf(self, file_path):
        """Extract text from PDF files with memory-efficient page-by-page processing"""
        try:
            print(f"  Opening PDF: {file_path}")
            reader = PdfReader(file_path)
            total_pages = len(reader.pages)
            print(f"  PDF has {total_pages} pages")
            
            # For very large PDFs, process page by page to avoid memory issues
            text_chunks = []
            max_pages_per_chunk = 2  # Process 2 pages at a time
            
            for i in range(0, total_pages, max_pages_per_chunk):
                chunk_text = ""
                end_page = min(i + max_pages_per_chunk, total_pages)
                
                print(f"  Processing PDF pages {i+1}-{end_page}/{total_pages}")
                
                # Process a few pages at a time
                for page_idx in range(i, end_page):
                    try:
                        page = reader.pages[page_idx]
                        page_text = page.extract_text()
                        if page_text:
                            chunk_text += page_text + "\n\n"
                        
                        # Release page object
                        del page
                    except Exception as e:
                        print(f"  Error extracting text from page {page_idx+1}: {str(e)}")
                
                # Add the chunk if it has content
                if chunk_text.strip():
                    text_chunks.append(chunk_text)
                
                # Force garbage collection after each chunk
                gc.collect()
                
                # For large PDFs, pause briefly to let memory stabilize
                if total_pages > 20 and i % 10 == 0:
                    print("  Pausing to stabilize memory...")
                    import time
                    time.sleep(0.5)
            
            # Combine chunks
            text = "\n".join(text_chunks)
            
            # Clean up
            del reader
            del text_chunks
            gc.collect()
            
            return text
        except Exception as e:
            print(f"Error extracting text from PDF {file_path}: {str(e)}")
            import traceback
            traceback.print_exc()
            return ""
    
    def extract_text_from_docx(self, file_path):
        """Extract text from DOCX files"""
        try:
            text = docx2txt.process(file_path)
            # Force garbage collection after processing
            gc.collect()
            return text
        except Exception as e:
            print(f"Error extracting text from DOCX {file_path}: {str(e)}")
            return ""
    
    def extract_text_from_txt(self, file_path):
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            print(f"Error extracting text from TXT {file_path}: {str(e)}")
            return ""
    
    def extract_text_from_csv(self, file_path):
        """Extract text from CSV files by converting to a string representation"""
        try:
            df = pd.read_csv(file_path)
            # Convert dataframe to string representation
            text = df.to_string()
            # Clean up
            del df
            gc.collect()
            return text
        except Exception as e:
            print(f"Error extracting text from CSV {file_path}: {str(e)}")
            return ""
            
    def extract_text_from_excel(self, file_path):
        """Extract text from Excel files by converting to a string representation"""
        try:
            # Read all sheets
            df_dict = pd.read_excel(file_path, sheet_name=None)
            text = ""
            
            # Process each sheet
            for sheet_name, df in df_dict.items():
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string() + "\n\n"
                
            # Clean up
            del df_dict
            gc.collect()
            return text
        except Exception as e:
            print(f"Error extracting text from Excel {file_path}: {str(e)}")
            return ""
    
    def extract_text_from_pptx(self, file_path):
        """Extract text from PowerPoint files"""
        try:
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            text = ""
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"  # Add extra line between slides
            
            # Clean up
            del presentation
            gc.collect()    
            return text
        except ImportError:
            print("python-pptx module not installed. Install it with: pip install python-pptx")
            return ""
        except Exception as e:
            print(f"Error extracting text from PPTX {file_path}: {str(e)}")
            return ""
    
    def extract_text(self, file_path):
        """Extract text from a file based on its extension"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Check if file is supported
        if file_extension not in self.supported_extensions:
            print(f"Unsupported file format: {file_extension}")
            return ""
            
        # Check file size
        file_size = os.path.getsize(file_path)
        if file_size > self.max_file_size_bytes:
            print(f"File too large ({file_size / (1024 * 1024):.2f} MB). Skipping.")
            return ""
        
        if file_extension == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_extension in ['.docx', '.doc']:
            return self.extract_text_from_docx(file_path)
        elif file_extension == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_extension == '.csv':
            return self.extract_text_from_csv(file_path)
        elif file_extension in ['.xlsx', '.xls']:
            return self.extract_text_from_excel(file_path)
        else:
            print(f"Unsupported file format: {file_extension}")
            return ""
    
    def clean_text(self, text):
        """Clean and normalize extracted text"""
        if not text:
            return ""
            
        # Replace multiple newlines with a single newline
        text = re.sub(r'\n+', '\n', text)
        # Replace multiple spaces with a single space
        text = re.sub(r'\s+', ' ', text)
        # Remove any special characters that might cause issues
        text = re.sub(r'[^\w\s\.\,\:\;\-\?\!]', ' ', text)
        
        return text.strip()
    
    def chunk_text(self, text, chunk_size=1000, overlap=200):
        """Split text into overlapping chunks for processing"""
        if not text:
            return []
            
        chunks = []
        text_length = len(text)
        start = 0
        
        # For very large texts, use a smaller chunk size to conserve memory
        if text_length > 100000:  # If text is over 100K characters
            print(f"  Very large text detected ({text_length} chars), using smaller chunks")
            chunk_size = min(chunk_size, 500)  # Cap chunk size at 500 chars for large documents
            overlap = min(overlap, 100)  # Reduce overlap for large documents
        
        chunk_count = 0
        while start < text_length:
            end = min(start + chunk_size, text_length)
            
            # If we're not at the end of the text, try to find a good break point
            if end < text_length:
                # Look for a period, newline, or space to break at
                for break_char in ['. ', '\n', ' ']:
                    pos = text.rfind(break_char, start, end)
                    if pos != -1:
                        end = pos + 1  # Include the break character
                        break
            
            # Extract the chunk and add it to our list
            chunk = text[start:end].strip()
            if chunk:  # Only add non-empty chunks
                chunks.append(chunk)
                chunk_count += 1
                
                # For large documents, perform garbage collection periodically
                if chunk_count % 10 == 0:
                    gc.collect()
            
            # Move to the next chunk, accounting for overlap
            start = end - overlap
            
            # Avoid getting stuck in an infinite loop if the overlap is too large
            if start <= 0 or start >= text_length:
                break
                
        return chunks
    
    def process_single_document(self, file_path, chunk_size=1000, overlap=200):
        """Process a single document and return text chunks with metadata"""
        if not os.path.exists(file_path):
            print(f"File does not exist: {file_path}")
            return []
            
        file_chunks = []
        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Skip unsupported file types
        if file_extension not in self.supported_extensions:
            print(f"Skipping unsupported file type: {file_path}")
            return []
            
        # Check file size
        file_size = os.path.getsize(file_path)
        if file_size > self.max_file_size_bytes:
            print(f"Skipping file that is too large ({file_size / (1024 * 1024):.2f} MB): {file_path}")
            return []
        
        print(f"Processing {file_name}...")
        
        # Extract text from the file
        text = self.extract_text(file_path)
        
        # Clean the extracted text
        cleaned_text = self.clean_text(text)
        
        # Free memory from the original text
        del text
        gc.collect()
        
        if cleaned_text:
            # Adjust chunk size for PDFs based on file size
            if file_extension == '.pdf' and file_size > 1024 * 1024:  # If PDF is over 1MB
                # Use smaller chunks for larger PDFs
                adjusted_chunk_size = min(chunk_size, 500)
                adjusted_overlap = min(overlap, 100)
                print(f"  Using smaller chunk size ({adjusted_chunk_size}) for large PDF")
                chunks = self.chunk_text(cleaned_text, chunk_size=adjusted_chunk_size, overlap=adjusted_overlap)
            else:
                # Use default chunking for other files
                chunks = self.chunk_text(cleaned_text, chunk_size=chunk_size, overlap=overlap)
            
            # Free memory from cleaned text
            del cleaned_text
            gc.collect()
            
            # Add metadata to each chunk
            for i, chunk in enumerate(chunks):
                chunk_data = {
                    "text": chunk,
                    "metadata": {
                        "source": file_name,
                        "chunk_id": i,
                        "total_chunks": len(chunks)
                    }
                }
                file_chunks.append(chunk_data)
                
                # For large documents, perform garbage collection periodically
                if i > 0 and i % 10 == 0:
                    gc.collect()
            
            print(f"  - Generated {len(chunks)} chunks")
            
            # Free memory from chunks list
            del chunks
            gc.collect()
        else:
            print(f"  - No text extracted")
        
        return file_chunks
    
    def process_documents(self):
        """Process all documents in the directory and return the text chunks with metadata"""
        all_chunks = []
        
        for root, _, files in os.walk(self.docs_directory):
            for file in files:
                file_path = os.path.join(root, file)
                relative_path = os.path.relpath(file_path, self.docs_directory)
                file_extension = os.path.splitext(file_path)[1].lower()
                
                # Skip unsupported file types
                if file_extension not in self.supported_extensions:
                    print(f"Skipping unsupported file type: {relative_path}")
                    continue
                    
                # Check file size
                file_size = os.path.getsize(file_path)
                if file_size > self.max_file_size_bytes:
                    print(f"Skipping file that is too large ({file_size / (1024 * 1024):.2f} MB): {relative_path}")
                    continue
                
                print(f"Processing {relative_path}...")
                
                # Process in a memory-efficient way
                file_chunks = self.process_single_document(file_path)
                all_chunks.extend(file_chunks)
                
                # Force garbage collection after each file
                del file_chunks
                gc.collect()
        
        return all_chunks

if __name__ == "__main__":
    processor = DocumentProcessor()
    chunks = processor.process_documents()
    print(f"Total chunks generated: {len(chunks)}") 